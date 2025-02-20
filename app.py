import pandas as pd
from pptx import Presentation
from pptx.util import Pt
import requests
from google import genai
import os
from dotenv import load_dotenv

load_dotenv()


client = genai.Client(api_key = os.getenv('API_KEY'))

# Função para gerar resumo usando a API do Gemini
def gerar_resumo(livro):
    response = client.models.generate_content(
        model="gemini-2.0-flash",
        contents=f"Faça um resumo de no máximo 445 caracteres sobre o livro: {livro}",
    )
    return response.text

def gerar_frase_motivacional(livros):
    temas = ", ".join(livros)  # Combina os títulos dos livros em uma única string
    response = client.models.generate_content(
        model="gemini-2.0-flash",
        contents=f"Crie apenas uma frase curta e motivacional de no máximo 100 caracteres que incentive a leitura (não quero sugestões, apenas retorne a frase sem mais nada além disso. não precisa deixar em negrito, então não coloque asteriscos '*'), baseada nos seguintes livros: {temas}",
    )
    return response.text

def baixar_imagem(url, caminho_local):
    resposta = requests.get(url)
    if resposta.status_code == 200:
        with open(caminho_local, "wb") as arquivo:
            arquivo.write(resposta.content)
    else:
        raise Exception(f"Erro ao baixar a imagem: {url}")


excel_path = 'livros_e_imagens.xlsx'

# Ler a aba "selecao" do arquivo Excel

df_selecao = pd.read_excel(excel_path, sheet_name='selecao')

# Ler os 3 primeiros livros e imagens

livros = df_selecao['livros'].head(3).tolist()
imagens = df_selecao['imagens'].head(3).tolist()

# Baixar as imagens e salvar localmente

caminhos_imagens = []
for i, url in enumerate(imagens):
    caminho_local = f"imagem_{i + 1}.jpg"  # Nome do arquivo local
    baixar_imagem(url, caminho_local)
    caminhos_imagens.append(caminho_local)


resumos = [gerar_resumo(livro) for livro in livros]


frase_motivacional = gerar_frase_motivacional(livros)


pptx_path = 'minha_apresentacao.pptx'


prs = Presentation(pptx_path)


def substituir_texto(slide, antigo_texto, novo_texto, tamanho_fonte=Pt(14)):
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                if antigo_texto in para.text:
                    para.text = para.text.replace(antigo_texto, novo_texto)
                    for run in para.runs:
                        run.font.size = tamanho_fonte


def substituir_imagem_por_nome(slide, nome_placeholder, nova_imagem):
    for shape in slide.shapes:
        if shape.name == nome_placeholder:
            left = shape.left
            top = shape.top
            width = shape.width
            height = shape.height
            slide.shapes.add_picture(nova_imagem, left, top, width, height)
            sp = shape._element
            sp.getparent().remove(sp)


for slide in prs.slides:
    substituir_texto(slide, 'texto1', resumos[0])
    substituir_texto(slide, 'texto2', resumos[1])
    substituir_texto(slide, 'texto3', resumos[2])
    substituir_texto(slide, 'texto4', frase_motivacional)  # Substituir o placeholder "texto4"
    substituir_imagem_por_nome(slide, 'imagem1', caminhos_imagens[0])
    substituir_imagem_por_nome(slide, 'imagem2', caminhos_imagens[1])
    substituir_imagem_por_nome(slide, 'imagem3', caminhos_imagens[2])


prs.save('apresentacao_modificada.pptx')


for caminho in caminhos_imagens:
    os.remove(caminho)

print("Apresentação atualizada com sucesso!")